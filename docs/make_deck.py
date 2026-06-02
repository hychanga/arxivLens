# -*- coding: utf-8 -*-
"""Generate a 3-slide arxivLens intro deck (architecture / features / data flow)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- palette (zinc + blue, matching the app) ----
INK     = RGBColor(0x18, 0x18, 0x1B)   # zinc-900
MUTED   = RGBColor(0x71, 0x71, 0x7A)   # zinc-500
ACCENT  = RGBColor(0x25, 0x63, 0xEB)   # blue-600
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LINE    = RGBColor(0xE4, 0xE4, 0xE7)   # zinc-200
BOX_FE  = RGBColor(0xEF, 0xF6, 0xFF)   # light blue
BOX_BE  = RGBColor(0xEC, 0xFD, 0xF5)   # light emerald
BOX_DB  = RGBColor(0xFE, 0xF3, 0xC7)   # light amber
BOX_GREY= RGBColor(0xF4, 0xF4, 0xF5)   # zinc-100
BD_FE   = RGBColor(0x93, 0xC5, 0xFD)
BD_BE   = RGBColor(0x6E, 0xE7, 0xB7)
BD_DB   = RGBColor(0xFC, 0xD3, 0x4D)

FONT = "Microsoft JhengHei"   # CJK-capable on Windows

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def _set_cjk(run, name=FONT):
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=2, line_spacing=1.05):
    """runs: list of paragraphs; each paragraph is a list of (txt, size, bold, color)."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (txt, size, bold, color) in para:
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            _set_cjk(r)
    return tb


def box(slide, x, y, w, h, fill, border, title=None, lines=None,
        title_size=13, line_size=10.5, rounded=True):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border
    shp.line.width = Pt(1)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(8); tf.margin_right = Pt(8)
    tf.margin_top = Pt(4); tf.margin_bottom = Pt(4)
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.CENTER
    if title:
        r = p0.add_run(); r.text = title
        r.font.size = Pt(title_size); r.font.bold = True; r.font.color.rgb = INK
        _set_cjk(r)
    for ln in (lines or []):
        p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER
        p.line_spacing = 1.0; p.space_after = Pt(0)
        r = p.add_run(); r.text = ln
        r.font.size = Pt(line_size); r.font.color.rgb = MUTED
        _set_cjk(r)
    return shp


def arrow(slide, x, y, w, h, char="→", size=20):
    text(slide, x, y, w, h, [[(char, size, True, ACCENT)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def header(slide, zh, en, n):
    # accent rule
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.55),
                                 Inches(0.12), Inches(0.62))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    bar.shadow.inherit = False
    text(slide, 0.85, 0.45, 10.5, 1.0,
         [[(zh + "  ", 26, True, INK), (en, 15, False, MUTED)]])
    text(slide, 12.0, 0.5, 0.9, 0.6, [[(f"{n}/3", 12, True, MUTED)]],
         align=PP_ALIGN.RIGHT)


def bg(slide):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = WHITE; r.line.fill.background()
    r.shadow.inherit = False
    slide.shapes._spTree.remove(r._element); slide.shapes._spTree.insert(2, r._element)


# ============================ Cover ============================
s = prs.slides.add_slide(BLANK); bg(s)
barc = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.95), Inches(2.5), Inches(0.18), Inches(2.0))
barc.fill.solid(); barc.fill.fore_color.rgb = ACCENT; barc.line.fill.background(); barc.shadow.inherit = False
text(s, 1.3, 2.35, 11.2, 1.2, [[("arxivLens", 54, True, INK)]])
text(s, 1.34, 3.6, 11.2, 0.7, [[("系統介紹   ", 24, True, ACCENT), ("System Introduction", 18, False, MUTED)]])
text(s, 1.34, 4.5, 11.2, 0.6, [[("arXiv / HBR 論文聚合閱讀器 — AI 摘要、關鍵字評分、多語翻譯", 15, False, INK)]])
dv = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.34), Inches(6.05), Inches(10.6), Inches(0.014))
dv.fill.solid(); dv.fill.fore_color.rgb = LINE; dv.line.fill.background(); dv.shadow.inherit = False
text(s, 1.34, 6.2, 11.2, 0.5,
     [[("Next.js 16 · React 19 · Spring Boot 4 · Java 25 · TiDB Cloud · Google Gemini 2.5 Flash", 12, False, MUTED)]])
text(s, 10.6, 0.7, 2.1, 0.4, [[("2026.06", 12, True, MUTED)]], align=PP_ALIGN.RIGHT)


# ============================ Slide 1 — Architecture ============================
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "系統架構", "System Architecture", 1)
text(s, 0.85, 1.35, 11.8, 0.5,
     [[("arxivLens — arXiv / HBR 論文聚合閱讀器，含 AI 摘要、關鍵字評分與多語翻譯", 13, False, MUTED)]])

row_y, bw, bh = 2.35, 2.55, 1.45
xs = [0.7, 3.65, 6.6, 9.55]
box(s, xs[0], row_y, bw, bh, BOX_GREY, LINE, "使用者瀏覽器", ["Browser", "桌機 / 行動裝置"])
box(s, xs[1], row_y, bw, bh, BOX_FE, BD_FE, "前端 Frontend", ["Vercel", "Next.js 16 · React 19", "TypeScript · Tailwind 4", "Zustand"])
box(s, xs[2], row_y, bw, bh, BOX_BE, BD_BE, "後端 Backend", ["Render (Docker)", "Spring Boot 4 · Java 25", "Spring Security · JWT"])
box(s, xs[3], row_y, bw, bh, BOX_DB, BD_DB, "資料庫 Database", ["TiDB Cloud Serverless", "(MySQL 相容)", "本地開發用 MySQL 8"])
for ax in (3.30, 6.25, 9.20):
    arrow(s, ax, row_y, 0.35, bh, "→")
# REST / JSON caption between FE and BE
text(s, 3.65, row_y + bh + 0.02, 5.9, 0.3,
     [[("HTTPS · REST / JSON · JWT", 9, False, MUTED)]], align=PP_ALIGN.CENTER)

# vertical arrow down from backend
arrow(s, xs[2] + bw/2 - 0.2, row_y + bh - 0.05, 0.4, 0.55, "↓", 18)

# external integrations row
ext_y = row_y + bh + 0.65
box(s, 0.7, ext_y, 11.4, 1.35, WHITE, LINE, "外部整合 External integrations", [])
chips = [
    ("arXiv Atom API", "論文同步"),
    ("HBR RSS", "文章來源"),
    ("Google Gemini 2.5 Flash", "摘要 / 翻譯"),
    ("Resend", "Email 通知"),
    ("GitHub Actions", "每 6 小時 cron"),
]
cw = 2.12; gap = 0.12; cx = 0.95
for name, sub in chips:
    box(s, cx, ext_y + 0.45, cw, 0.78, BOX_GREY, LINE, name, [sub],
        title_size=10.5, line_size=9)
    cx += cw + gap

text(s, 0.7, 6.95, 11.8, 0.4,
     [[("部署：Vercel（前端）· Render（後端）· TiDB Cloud（資料庫）— 全部免費方案", 10.5, False, MUTED)]])


# ============================ Slide 2 — Features ============================
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "系統特色", "Key Features", 2)

groups = [
    ("內容與閱讀", "Feed & Reading", [
        "Latest 聚合 arXiv + HBR，依主題 / 天數過濾",
        "關鍵字相關性評分，重點論文一眼可見",
        "交叉分類標示（cs.CV 論文掛 cs.AI 也會顯示 +cs.AI）",
        "Library 本地快取 PDF，內建檢視器（讀取時 loading 游標）",
    ]),
    ("AI 與多語", "AI & i18n", [
        "Gemini 2.5 Flash：依需求生成 AI 摘要",
        "標題 / 摘要翻譯 zh-TW · zh-CN · ja · de",
        "每篇翻譯快取於 DB，整頁批次查詢（省連線）",
    ]),
    ("趨勢分析", "Trends", [
        "12 個月主題堆疊長條圖",
        "採 arXiv 官方總量（opensearch totalResults）",
    ]),
    ("帳號與安全", "Auth & Security", [
        "Email / 密碼登入 + JWT",
        "兩階段驗證 2FA（TOTP）",
        "Sign in with Google / Apple、密碼重設",
    ]),
    ("維運自動化", "Ops", [
        "Admin：來源 / 主題 / 設定、深度 resync、backfill",
        "每 6 小時自動同步 + 郵件通知",
        "外部 cron 觸發（克服 Render 免費方案休眠）",
    ]),
]
# two columns
col_x = [0.7, 6.95]; col_w = 5.9
col_y = [1.55, 1.55]
order = [0, 1, 2, 3, 4]
# place: left col groups 0,1,2 ; right col 3,4
layout = {0: 0, 1: 0, 2: 0, 3: 1, 4: 1}
for gi in order:
    col = layout[gi]
    zh, en, items = groups[gi]
    y = col_y[col]
    # group header
    text(s, col_x[col], y, col_w, 0.4,
         [[("▍" + zh + "  ", 15, True, ACCENT), (en, 11, False, MUTED)]])
    y += 0.45
    runs = [[("•  ", 12, False, ACCENT), (it, 12, False, INK)] for it in items]
    tb = text(s, col_x[col], y, col_w, 0.32 * len(items) + 0.1, runs,
              line_spacing=1.1, space_after=3)
    y += 0.30 * len(items) + 0.30
    col_y[col] = y


# ============================ Slide 3 — Data Flow ============================
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "資料流", "Data Flow", 3)

# Flow A — Sync (every 6h)
text(s, 0.7, 1.45, 11.8, 0.4,
     [[("①  同步流程 Sync — 每 6 小時自動", 15, True, INK)]])
ay, abw, abh = 1.95, 1.95, 1.15
ax = 0.7
sync_steps = [
    ("GitHub Actions", "cron 0 */6", BOX_GREY, LINE),
    ("/api/cron/arxiv-sync", "token 驗證", BOX_BE, BD_BE),
    ("ArxivSyncService", "分頁・退避重試", BOX_BE, BD_BE),
    ("arXiv Atom API", "lastUpdatedDate", BOX_GREY, LINE),
    ("逐頁交易 upsert", "categories 回填", BOX_BE, BD_BE),
    ("TiDB", "papers", BOX_DB, BD_DB),
]
n = len(sync_steps); abw = (11.4 - 0.30 * (n - 1)) / n
ax = 0.7
for i, (t1, t2, f, b) in enumerate(sync_steps):
    box(s, ax, ay, abw, abh, f, b, t1, [t2], title_size=10.5, line_size=9)
    if i < n - 1:
        arrow(s, ax + abw - 0.02, ay, 0.34, abh, "→", 16)
    ax += abw + 0.30
# tail: notify email
arrow(s, 0.7 + abw/2 - 0.2, ay + abh - 0.05, 0.4, 0.5, "↓", 15)
box(s, 0.7, ay + abh + 0.45, 4.2, 0.7, BOX_GREY, LINE,
    "Resend 寄出摘要信", ["新增 / 跳過 / 錯誤統計 → hychanga@gmail.com"],
    title_size=10.5, line_size=9)

# Flow B — Read (on demand)
by = 4.75
text(s, 0.7, by, 11.8, 0.4,
     [[("②  閱讀流程 Read — 即時依需求", 15, True, INK)]])
ry = by + 0.5
read_steps = [
    ("使用者 / Browser", "Next.js Feed", BOX_FE, BD_FE),
    ("GET /api/papers", "主題 / 天數 / 分頁", BOX_BE, BD_BE),
    ("TiDB", "查詢已同步論文", BOX_DB, BD_DB),
    ("論文卡片", "評分 · 翻譯 · 收藏", BOX_FE, BD_FE),
]
n2 = len(read_steps); rbw = (11.4 - 0.30 * (n2 - 1)) / n2
rx = 0.7
for i, (t1, t2, f, b) in enumerate(read_steps):
    box(s, rx, ry, rbw, 1.0, f, b, t1, [t2], title_size=10.5, line_size=9)
    if i < n2 - 1:
        arrow(s, rx + rbw - 0.02, ry, 0.34, 1.0, "→", 16)
    rx += rbw + 0.30
text(s, 0.7, ry + 1.1, 11.8, 0.4,
     [[("翻譯 / 摘要：批次 GET /api/papers/translations → 命中 DB 快取，否則呼叫 Gemini 生成後寫回快取", 10.5, False, MUTED)]])

out = r"D:\greg\project\arxivlens\docs\arxivLens-系統介紹.pptx"
prs.save(out)
print("saved:", out)
